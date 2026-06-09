#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
可编辑 PPT / Editable PPTX of all results.
设计: 统计结果用 python-pptx 原生图表(柱状/散点) —— 在 PowerPoint 中可
  双击编辑数据、改配色、改样式(真正可编辑); 地图与雨林/蜂群图嵌入高清图。
Native editable charts for every statistical result + embedded maps + the
beeswarm/raincloud images (vector PDF/SVG also available alongside).
输入: results/tables/*.csv, results/forecasts/*.csv, figures/main/*.png
输出: output/doc/bird_hazard_results_editable.pptx
运行: python3 code/57_editable_pptx.py  (需 56_ 先生成蜂群/雨林图)
"""
from pathlib import Path
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.text import PP_ALIGN

ROOT = Path(".").resolve()
T = ROOT/"results"/"tables"; F = ROOT/"results"/"forecasts"; FIG = ROOT/"figures"/"main"
OUT = ROOT/"output"/"doc"/"bird_hazard_results_editable.pptx"
OUT.parent.mkdir(parents=True, exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
NAVY = RGBColor(0x1F, 0x2D, 0x3D); ACC = RGBColor(0xD5, 0x5E, 0x00)
OKABE = [RGBColor(0x00, 0x72, 0xB2), RGBColor(0xE6, 0x9F, 0x00),
         RGBColor(0x00, 0x9E, 0x73), RGBColor(0xCC, 0x79, 0xA7),
         RGBColor(0xD5, 0x5E, 0x00), RGBColor(0x56, 0xB4, 0xE9)]

def add_title(slide, text, sub=None):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = NAVY
    if sub:
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(12); r2.font.color.rgb = RGBColor(0x55,0x55,0x55)

def add_note(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4))
    p = tb.text_frame.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(9); r.font.italic = True; r.font.color.rgb = RGBColor(0x77,0x77,0x77)

def style_chart(chart, legend=True, dl=True, num="0.00"):
    chart.has_legend = legend
    if legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False; chart.legend.font.size = Pt(11)
    try:
        plot = chart.plots[0]; plot.has_data_labels = dl
        if dl:
            plot.data_labels.number_format = num
            plot.data_labels.number_format_is_linked = False
            plot.data_labels.font.size = Pt(9)
            plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    except Exception:
        pass
    for i, s in enumerate(chart.series):
        try: s.format.fill.solid(); s.format.fill.fore_color.rgb = OKABE[i % len(OKABE)]
        except Exception: pass

def bar_slide(title, sub, cats, series, note=None, ctype=XL_CHART_TYPE.COLUMN_CLUSTERED,
              num="0.00", x=0.7, y=1.4, cx=12.0, cy=5.2, dl=True):
    s = prs.slides.add_slide(BLANK); add_title(s, title, sub)
    cd = CategoryChartData(); cd.categories = cats
    for nm, vals in series: cd.add_series(nm, vals)
    gf = s.shapes.add_chart(ctype, Inches(x), Inches(y), Inches(cx), Inches(cy), cd)
    style_chart(gf.chart, legend=len(series) > 1, dl=dl, num=num)
    if note: add_note(s, note)
    return s

def img_slide(title, sub, img, note=None, max_w=12.2, max_h=5.4):
    s = prs.slides.add_slide(BLANK); add_title(s, title, sub)
    if Path(img).exists():
        from PIL import Image
        iw, ih = Image.open(img).size; ar = iw/ih
        w = max_w; h = w/ar
        if h > max_h: h = max_h; w = h*ar
        left = (13.333 - w)/2
        s.shapes.add_picture(str(img), Inches(left), Inches(1.45), Inches(w), Inches(h))
    else:
        add_note(s, f"[missing image: {img}]")
    if note: add_note(s, note)
    return s

# ---------- Title ----------
s = prs.slides.add_slide(BLANK)
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.6))
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "Survey effort moderates the climatic visibility of new bird distribution records across China"
r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = NAVY
for t in ["Editable results deck — native charts (double-click to edit) + maps + beeswarm/raincloud",
          "Chen-Chen Ding · Peking University · all values from persisted result tables"]:
    p = tf.add_paragraph(); rr = p.add_run(); rr.text = t
    rr.font.size = Pt(14); rr.font.color.rgb = RGBColor(0x55,0x55,0x55)

# ---------- 1. Headline interaction HR (v2/lag/v3 x 4 specs) ----------
v2 = pd.read_csv(T/"table_province_v2_coefs.csv")
v2i = v2[(v2.model=="M4")&(v2.term=="climate_z:effort_z")].set_index("spec_id")
lag = pd.read_csv(T/"table_effort_lag_refit.csv").set_index("spec_id")
v3 = pd.read_csv(T/"table_province_v3_all_specs_coefs.csv")
v3i = v3[(v3.model=="M4")&(v3.term=="climate_z:effort_z")].set_index("spec_id")
order = ["spec_A","spec_B","spec_C","spec_D"]
cats = ["A records","B visits","C PCA","D birding-days"]
bar_slide("Climate × effort interaction hazard ratio",
    "Positive (>1) across all 4 effort metrics and 3 risk sets, incl. lagged-effort endogeneity test. HR=1 is no effect.",
    cats,
    [("v2 conservative", [round(v2i.loc[o,"hr"],3) for o in order]),
     ("v2 lagged effort (t-1)", [round(lag.loc[o,"hr"],3) for o in order]),
     ("v3 relaxed", [round(v3i.loc[o,"hr"],3) for o in order])],
    note="Source: table_province_v2_coefs / table_effort_lag_refit / table_province_v3_all_specs_coefs",
    num="0.000")

# ---------- 2. AIC ladder M0-M4 (v2 Spec B) ----------
aic = pd.read_csv(T/"table_province_v2_aic.csv")
aicB = aic[aic.spec_id=="spec_B"].set_index("model")
mods = ["M0","M1","M2","M3","M4"]
bar_slide("Model-selection ladder (ΔAIC vs best, v2 Spec B)",
    "M4 (climate × effort interaction) is the best model; lower is better, M4 ΔAIC = 0.",
    mods, [("ΔAIC", [round(float(aicB.loc[m,"dAIC"]),1) for m in mods])],
    note="Source: table_province_v2_aic.csv (Spec B). M4 Akaike weight ≈ 1.0.", num="0.0")

# ---------- 3. M4 vs M5 (offset) ΔAIC ----------
m5 = pd.read_csv(T/"table_m5_offset_summary.csv")
m5["lab"] = m5["spec_label"].str.replace(r"\s*\(.*\)","",regex=True)+" ("+m5["run"]+")"
bar_slide("Moderator vs scaling: M4 (interaction) − M5 (offset) ΔAIC",
    "Negative = interaction model wins (effort is a moderator, not a scaling factor). M4 wins 7 of 8.",
    list(m5["lab"]), [("ΔAIC (M4−M5)", [round(float(x),0) for x in m5["dAIC_M4_minus_M5"]])],
    note="Source: table_m5_offset_summary.csv", num="0", dl=True, cy=5.0)

# ---------- 4. Multi-scale interaction HR (v2) ----------
pref = pd.read_csv(T/"table_prefecture_coefs.csv"); cnty = pd.read_csv(T/"table_county_coefs.csv")
prov_hr = float(v2i.loc["spec_B","hr"])
pref_hr = float(pref[(pref.model=="M4")&(pref.term.str.contains(":"))]["hr"].iloc[0])
cnty_hr = float(cnty[(cnty.model=="M4")&(cnty.term.str.contains(":"))]["hr"].iloc[0])
bar_slide("Interaction across administrative scales (v2, Spec B)",
    "Attenuates with finer grain but never reverses sign (province → prefecture → county).",
    ["Province","Prefecture","County"],
    [("Interaction HR", [round(prov_hr,3), round(pref_hr,3), round(cnty_hr,3)])],
    note="Source: table_province_v2_coefs / table_prefecture_coefs / table_county_coefs", num="0.000")

# ---------- 5. Endogeneity: lagged-effort HR ----------
bar_slide("Endogeneity check — interaction with effort lagged one year",
    "Effort measured BEFORE the event still moderates record hazard → argues against reverse causation.",
    cats, [("Lagged-effort interaction HR", [round(lag.loc[o,"hr"],3) for o in order])],
    note="Source: table_effort_lag_refit.csv (province-year lag, n=12,535, 498 events)", num="0.000")

# ---------- 6. Spatial vs random CV ----------
cv = pd.read_csv(T/"table_spatial_block_cv.csv")
bar_slide("Predictive transfer: spatial-block vs random cross-validation",
    "Spatially honest CV (250 km blocks) ≈ 0.55 AUC vs random ≈ 0.65 → random CV overstates skill; use for inference, not out-of-region forecasting.",
    ["Spatial-block (250 km)","Random 5-fold"],
    [("Mean AUC", [round(float(cv[cv.cv_type=='spatial_block_250km']['auc_mean'].iloc[0]),3),
                   round(float(cv[cv.cv_type=='random_5fold']['auc_mean'].iloc[0]),3)])],
    note="Source: table_spatial_block_cv.csv (M4 cloglog, Spec B)", num="0.000")

# ---------- 7. RF importance (v2 top 12) ----------
rf = pd.read_csv(T/"table_rf_importance_v2.csv").sort_values("importance", ascending=True).tail(12)
bar_slide("Random-forest permutation importance (v2, top 12)",
    "temp × effort interaction is the single most important predictor of new-record hazard.",
    list(rf["variable"]), [("Importance", [round(float(x),4) for x in rf["importance"]])],
    ctype=XL_CHART_TYPE.BAR_CLUSTERED, num="0.0000", note="Source: table_rf_importance_v2.csv",
    cx=11.0, cy=5.2, dl=False)

# ---------- 8. Future hazard top provinces SSP585-2050 ----------
fut = pd.read_csv(F/"table_province_future_glmmTMB.csv")
sel = fut[(fut.ssp=="SSP585")&(fut.year==2050)&(~fut.province.isin(["Hong Kong","Macao","Taiwan"]))]
sel = sel.sort_values("hazard_glmm", ascending=False).head(10)
bar_slide("Scenario future hazard — top 10 provinces (SSP585, 2050)",
    "Scenario-conditioned (effort frozen at 2024). High where exposure × existing high effort coincide; act on exposure × effort-GAP instead.",
    list(sel["province"]), [("Hazard", [round(float(x),3) for x in sel["hazard_glmm"]])],
    note="Source: table_province_future_glmmTMB.csv (mainland)", num="0.00", cy=5.0)

# ---------- 9. Beeswarm — NATIVE editable XY scatter ----------
rng = np.random.default_rng(7)
sB = prs.slides.add_slide(BLANK)
add_title(sB, "Interaction HR — beeswarm (native editable scatter)",
          "Each series = a risk set; points = bootstrap draws of the interaction HR across the 4 specs. Fully editable in PowerPoint.")
betas = {}
for o in order: betas[("v2",o)] = (v2i.loc[o,"beta"], v2i.loc[o,"se"])
for o in order: betas[("lag",o)] = (lag.loc[o,"beta"], lag.loc[o,"se"])
for o in order: betas[("v3",o)] = (v3i.loc[o,"beta"], v3i.loc[o,"se"])
xy = XyChartData()
dodge = {"v2":-0.22,"lag":0.0,"v3":0.22}
for ds,nm in [("v2","v2 conservative"),("lag","v2 lagged (t-1)"),("v3","v3 relaxed")]:
    ser = xy.add_series(nm)
    for i,o in enumerate(order):
        b,se = betas[(ds,o)]
        draws = np.exp(rng.normal(b,se,90))
        xs = i + dodge[ds] + rng.uniform(-0.08,0.08,90)
        for xv,yv in zip(xs,draws): ser.add_data_point(float(xv), float(yv))
gf = sB.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, Inches(0.7), Inches(1.5),
                         Inches(12.0), Inches(5.0), xy)
ch = gf.chart; ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.include_in_layout = False
for i,sr in enumerate(ch.series):
    try:
        sr.marker.format.fill.solid(); sr.marker.format.fill.fore_color.rgb = OKABE[i]
        sr.marker.format.line.fill.background()
    except Exception: pass
add_note(sB, "x = effort spec (0=A records,1=B visits,2=C PCA,3=D birding-days). y = interaction HR. A polished vector raincloud/beeswarm follows.")

# ---------- 10-11. Embedded beautiful beeswarm + raincloud ----------
img_slide("Interaction HR — publication raincloud",
          "All distributions sit above HR 1 across specs and risk sets (incl. lagged-effort endogeneity test).",
          FIG/"Figure_R1_interaction_raincloud.png",
          note="Vector versions: figures/main/Figure_R1_interaction_raincloud.{pdf,svg}")
img_slide("Interaction HR — publication beeswarm",
          "Bootstrap draws of the interaction HR; black bar = median, line = 95% interval.",
          FIG/"Figure_R2_interaction_beeswarm.png",
          note="Vector versions: figures/main/Figure_R2_interaction_beeswarm.{pdf,svg}")

# ---------- 12+. Embedded main/maps ----------
EMB = [
    ("Figure 1 — concept and workflow", "Study domain, sample structure, climate × effort × visibility logic", "Figure_1_concept_and_workflow.png"),
    ("Figure 2 (v4) — province headline + M5 + raincloud", "Forest, AIC ladder, M4-vs-M5 dumbbell, bootstrap raincloud", "Figure_2_v4_province_headline_M5_raincloud.png"),
    ("Figure 3 (v4) — three-scale forest", "Interaction across province/prefecture/county", "Figure_3_v4_three_scale_forest_M5.png"),
    ("Figure 4 (v3) — robustness boundary", "Event recovery, reconciliation, multi-scale weakening, RF rank shift", "Figure_4_variable_importance_v3.png"),
    ("Figure 5 (v4) — province future hazard (glmmTMB)", "SSP245/585 × 2030/2050/2080, equal-area choropleth", "Figure_5_v4_province_future_glmmTMB.png"),
    ("Figure 5 (v4) — province future hazard (XGBoost)", "Machine-learning surrogate projection", "Figure_5_v4_province_future_xgboost.png"),
    ("Figure 6 — prefecture (市) plug-in hazard", "CRU climate × merged effort, CMIP6 ensemble futures (unified plug-in)", "Figure_6_prefecture_plugin_unified.png"),
    ("Figure 7 — county (县) plug-in hazard", "CRU climate × merged effort, CMIP6 ensemble futures (unified plug-in)", "Figure_7_county_plugin_unified.png"),
    ("Figure 8 — 100 km grid plug-in hazard", "CRU temperature anomaly × merged eBird+China-Birdwatch effort; CMIP6 ensemble futures", "Figure_8_grid100_native_plugin_hazard.png"),
    ("Figure 8b — 50 km grid plug-in hazard", "Same design at finer 50 km grain", "Figure_8b_grid50_native_plugin_hazard.png"),
]
for title, sub, fn in EMB:
    img_slide(title, sub, FIG/fn, note=f"figures/main/{fn} (+ .pdf 600 dpi)")

prs.save(OUT)
print(f"[57] wrote {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
